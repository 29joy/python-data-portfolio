import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

# 1. Read Excel
df = pd.read_excel("sales.xlsx")

# Data Summary (Monthly Product Sales)
pivot = df.pivot_table(index="Month", columns="Product", values="Sales", aggfunc="sum")

# 2. Create PPT
prs = Presentation()

# ========== Slide 1: Table ==========
slide_layout = prs.slide_layouts[5]  # blank page
slide = prs.slides.add_slide(slide_layout)

# Add title
title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
title_tf = title_shape.text_frame
title_tf.text = "Sales Data (Raw)"
title_tf.paragraphs[0].font.size = Pt(24)

# Add table
rows, cols = df.shape[0] + 1, df.shape[1]
table = slide.shapes.add_table(
    rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4)
).table

# Header
for j, col in enumerate(df.columns):
    table.cell(0, j).text = col

# Padding data
for i in range(df.shape[0]):
    for j in range(df.shape[1]):
        table.cell(i + 1, j).text = str(df.iloc[i, j])

# ========== Slide 2: Bar Chart ==========
slide = prs.slides.add_slide(slide_layout)

# Add title
title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
title_shape.text_frame.text = "Monthly Sales by Product"

# Chart data
chart_data = CategoryChartData()
chart_data.categories = pivot.index
for col in pivot.columns:
    chart_data.add_series(col, pivot[col].values)

# Add bar chart
x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4.5)
slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

# 3. save to PPT
prs.save("sales_report.pptx")
print("✅ PowerPoint saved as sales_report.pptx")
